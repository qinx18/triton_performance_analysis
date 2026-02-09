#!/usr/bin/env python3
"""Generate PowerPoint slides for the Analysis-Guided LLM Tritonization Pipeline."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MED_BLUE = RGBColor(0x2E, 0x5C, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
ACCENT_ORANGE = RGBColor(0xE8, 0x7D, 0x2F)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
MED_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
CODE_BG = RGBColor(0xF0, 0xF3, 0xF4)
TABLE_HEADER_BG = RGBColor(0x2E, 0x5C, 0x8A)

FONT_BODY = "Calibri"
FONT_CODE = "Courier New"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bold_prefix=False, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for idx, item in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.line_spacing = Pt(font_size * line_spacing)
        p.space_after = Pt(4)

        if bold_prefix and ": " in item:
            prefix, rest = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = prefix + ": "
            run1.font.name = FONT_BODY
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run2 = p.add_run()
            run2.text = rest
            run2.font.name = FONT_BODY
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = item
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return txBox


def add_code_box(slide, left, top, width, height, code_text, font_size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = MED_GRAY
    shape.line.width = Pt(0.5)
    # Smaller corner radius
    shape.adjustments[0] = 0.02

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    lines = code_text.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT_CODE
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.line_spacing = Pt(font_size * 1.3)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=MED_BLUE, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_flow_box(slide, left, top, width, height, text, fill_color=LIGHT_BLUE,
                 text_color=WHITE, font_size=11, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.15

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_BODY
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_section_title(slide, text, top=Inches(0.3)):
    add_textbox(slide, Inches(0.6), top, Inches(8.4), Inches(0.6),
                text, font_size=28, bold=True, color=DARK_BLUE)


def add_horizontal_arrow(slide, x, y, length=Inches(0.3)):
    """Draw a right-pointing arrow using a triangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, length, Inches(0.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MED_BLUE
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, x, y, length=Inches(0.25)):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, Inches(0.25), length)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MED_BLUE
    shape.line.fill.background()
    return shape


def add_simple_table(slide, left, top, width, col_widths, rows, header=True):
    """Add a table. rows[0] = header if header=True. col_widths in inches."""
    tbl = slide.shapes.add_table(len(rows), len(col_widths), left, top, width,
                                 Inches(0.35 * len(rows))).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(13)

            if header and ri == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            else:
                p.font.color.rgb = DARK_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else LIGHT_GRAY
    return tbl


# ════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BLUE)

    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(8.4), Inches(1.2),
                "Analysis-Guided LLM\nTritonization Pipeline",
                font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(8.4), Inches(0.7),
                "PET + LLVM Static Analysis  ->  LLM Prompt  ->  GPU Triton Kernels",
                font_size=18, color=NEAR_WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(4.1), Inches(8.4), Inches(0.5),
                "Polybench/C 4.2.1  |  30 HPC Kernels  |  16 Analysis Modules",
                font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


def slide_02_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Pipeline Architecture")

    # Row 1: main flow
    y1 = Inches(1.2)
    bw, bh = Inches(1.2), Inches(0.45)
    gap = Inches(0.15)

    boxes = [
        ("C Kernel", DARK_BLUE),
        ("Extract\nScop", MED_BLUE),
        ("PET / LLVM\nAnalysis", LIGHT_BLUE),
        ("Build\nPrompt", MED_BLUE),
        ("Claude\nSonnet", ACCENT_ORANGE),
        ("Triton\nCode", ACCENT_GREEN),
        ("Test vs\nC Ref", MED_BLUE),
    ]

    x = Inches(0.25)
    positions = []
    for label, color in boxes:
        add_flow_box(slide, x, y1, bw, bh, label, fill_color=color, font_size=9)
        positions.append(x)
        x_arrow = x + bw + gap * 0.1
        if label != "Test vs\nC Ref":
            add_horizontal_arrow(slide, x_arrow, y1 + bh / 2 - Inches(0.1),
                                 length=gap * 0.8)
        x = x_arrow + gap * 0.8

    # Retry loop arrow (text annotation)
    add_textbox(slide, Inches(7.2), Inches(1.8), Inches(2.5), Inches(0.35),
                "Retry loop (up to 5x)",
                font_size=10, color=ACCENT_RED, bold=True)

    # Stats bar
    stats_y = Inches(2.3)
    stats = [
        ("30 Kernels", DARK_BLUE),
        ("16 Analysis Modules", MED_BLUE),
        ("Up to 5 Attempts", ACCENT_ORANGE),
        ("25/30 Pass (83%)", ACCENT_GREEN),
    ]
    x = Inches(0.5)
    for label, color in stats:
        add_flow_box(slide, x, stats_y, Inches(2.0), Inches(0.4), label,
                     fill_color=color, font_size=11)
        x += Inches(2.3)

    # Analysis modules breakdown
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(9.0), Inches(0.4),
                "16 Analysis Modules (PET + LLVM fallback):",
                font_size=14, bold=True, color=DARK_BLUE)

    modules_text = [
        "WAR Dependencies  |  Parallel Dimensions  |  Reduction Detection  |  Scalar Expansion",
        "Stream Patterns  |  Aliasing  |  Overwrites  |  Indirect Addressing  |  Goto Detection",
        "Loop Distribution  |  Loop Interchange  |  Unrolling  |  Reordering  |  Crossing  |  Convolution  |  Early Exit",
    ]
    add_bullet_list(slide, Inches(0.5), Inches(3.6), Inches(9.0), Inches(2.0),
                    modules_text, font_size=12, color=DARK_TEXT)


def slide_03_pet(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "PET -- Polyhedral Extraction Tool")

    add_bullet_list(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(1.4), [
        "What: Parses #pragma scop C regions into a polyhedral model (ISL)",
        "Output: Iteration domains, access relations, schedules",
        "Why polyhedral? Mathematically precise dependency analysis",
    ], font_size=15, bold_prefix=True)

    # GEMM Example header
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(9.0), Inches(0.4),
                "GEMM Example:", font_size=16, bold=True, color=MED_BLUE)

    # C code
    add_textbox(slide, Inches(0.5), Inches(2.9), Inches(4.0), Inches(0.3),
                "C Kernel:", font_size=12, bold=True, color=DARK_TEXT)
    gemm_c = """\
for (i = 0; i < NI; i++) {
  for (j = 0; j < NJ; j++)
    C[i][j] *= beta;
  for (k = 0; k < NK; k++)
    for (j = 0; j < NJ; j++)
      C[i][j] += alpha*A[i][k]*B[k][j];
}"""
    add_code_box(slide, Inches(0.5), Inches(3.2), Inches(4.2), Inches(1.6), gemm_c, font_size=10)

    # PET output
    add_textbox(slide, Inches(5.0), Inches(2.9), Inches(4.8), Inches(0.3),
                "PET Polyhedral Output:", font_size=12, bold=True, color=DARK_TEXT)
    pet_out = """\
Domain S_0: { S_0[i,j] :
  0 <= i < 60 and 0 <= j < 70 }

Domain S_1: { S_1[i,k,j] :
  0 <= i<60, 0<=k<80, 0<=j<70 }

Write: S_0[i,j] -> C[i,j]
Read:  S_1[i,k,j] -> A[i,k]
Read:  S_1[i,k,j] -> B[k,j]
Write: S_1[i,k,j] -> C[i,j]"""
    add_code_box(slide, Inches(5.0), Inches(3.2), Inches(4.5), Inches(1.6), pet_out, font_size=9)


def slide_04_llvm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "LLVM Analysis Infrastructure")

    # 4 layers diagram
    layers = [
        ("Clang AST (JSON)", "Structure: loops, branches, array refs, function signatures",
         DARK_BLUE),
        ("LLVM IR", "Normalized representation: SSA form, optimized at -O1",
         MED_BLUE),
        ("DependenceAnalysis", "Flow / Anti / Output deps:  opt -passes='print<da>'",
         LIGHT_BLUE),
        ("ScalarEvolution (SCEV)", "Loop bounds, strides, trip counts, induction variables",
         ACCENT_GREEN),
    ]

    y = Inches(1.2)
    for label, desc, color in layers:
        add_flow_box(slide, Inches(0.5), y, Inches(2.5), Inches(0.5), label,
                     fill_color=color, font_size=12)
        add_textbox(slide, Inches(3.2), y + Inches(0.05), Inches(6.3), Inches(0.5),
                    desc, font_size=13, color=DARK_TEXT)
        y += Inches(0.7)

    # Advantage
    y += Inches(0.2)
    add_textbox(slide, Inches(0.5), y, Inches(9.0), Inches(0.4),
                "Advantage over PET alone:",
                font_size=16, bold=True, color=MED_BLUE)
    add_bullet_list(slide, Inches(0.5), y + Inches(0.4), Inches(9.0), Inches(1.5), [
        "Works on any code Clang can compile (no polyhedral restriction)",
        "Handles non-affine, pointer arithmetic, complex control flow",
        "PET requires #pragma scop + affine loops; LLVM has no such constraints",
    ], font_size=14)


def slide_05_fallback(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "PET <-> LLVM Fallback Strategy")

    # Code showing the pattern
    code = """\
def try_with_llvm_fallback(pet_func, llvm_fallback_func, *args):
    \"\"\"Try PET first, fall back to LLVM on failure.\"\"\"
    try:
        result = pet_func(*args)
        if result is not None:
            return result           # PET succeeded
    except Exception:
        pass
    try:
        return llvm_fallback_func(*args)  # LLVM fallback
    except Exception:
        return None                 # Both failed gracefully"""
    add_code_box(slide, Inches(0.5), Inches(1.1), Inches(6.5), Inches(2.3), code, font_size=11)

    # Robustness results
    add_textbox(slide, Inches(7.3), Inches(1.1), Inches(2.5), Inches(0.4),
                "Robustness", font_size=18, bold=True, color=MED_BLUE)

    robustness_items = [
        "0 crashes across 480 combos",
        "(16 modules x 30 kernels)",
        "",
        "8 modules: 100% pass",
        "ScalarExp: 90% pass",
        "(3 kernels empty output)",
        "",
        "Same dict format from",
        "both PET and LLVM paths",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(1.6), Inches(2.5), Inches(2.5),
                    robustness_items, font_size=12, color=DARK_TEXT)

    # Flow diagram
    y_flow = Inches(3.8)
    add_flow_box(slide, Inches(0.8), y_flow, Inches(1.8), Inches(0.5),
                 "C Kernel File", fill_color=DARK_BLUE, font_size=11)
    add_horizontal_arrow(slide, Inches(2.7), y_flow + Inches(0.15))

    add_flow_box(slide, Inches(3.2), y_flow, Inches(1.8), Inches(0.5),
                 "PET Analysis", fill_color=MED_BLUE, font_size=11)

    # Success path
    add_horizontal_arrow(slide, Inches(5.1), y_flow + Inches(0.15))
    add_flow_box(slide, Inches(5.6), y_flow, Inches(1.6), Inches(0.5),
                 "Result Dict", fill_color=ACCENT_GREEN, font_size=11)

    # Fail path
    add_down_arrow(slide, Inches(3.95), y_flow + Inches(0.55))
    add_flow_box(slide, Inches(3.2), y_flow + Inches(0.95), Inches(1.8), Inches(0.5),
                 "LLVM Fallback", fill_color=ACCENT_ORANGE, font_size=11)
    add_horizontal_arrow(slide, Inches(5.1), y_flow + Inches(1.1))
    add_flow_box(slide, Inches(5.6), y_flow + Inches(0.95), Inches(1.6), Inches(0.5),
                 "Result Dict", fill_color=ACCENT_GREEN, font_size=11)

    add_textbox(slide, Inches(3.3), y_flow + Inches(0.5), Inches(1.5), Inches(0.35),
                "(fail/None)", font_size=9, color=ACCENT_RED, bold=True)


def slide_06_war(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "WAR Dependency Analysis (Example)")

    # S115 kernel
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.3),
                "Kernel s115:  a[i] -= aa[j][i] * a[j]",
                font_size=16, bold=True, color=MED_BLUE)

    code_s115 = """\
for (i = 0; i < LEN_1D; i++)
  for (j = 0; j < i; j++)
    a[i] -= aa[j][i] * a[j];   // reads a[j], writes a[i], j < i"""
    add_code_box(slide, Inches(0.5), Inches(1.4), Inches(5.5), Inches(1.0), code_s115, font_size=11)

    # PET detection
    add_textbox(slide, Inches(0.5), Inches(2.6), Inches(9.0), Inches(0.3),
                "PET Detects:", font_size=15, bold=True, color=DARK_BLUE)

    add_bullet_list(slide, Inches(0.5), Inches(2.95), Inches(9.0), Inches(1.5), [
        "Read access: a[j] (inner index j < i)",
        "Write access: a[i] (outer index i)",
        "Read index != Write index on same array -> WAR conflict",
        "Concurrent threads: thread_j may overwrite a[j] before thread_i reads it",
    ], font_size=13)

    # Analysis output
    add_textbox(slide, Inches(0.5), Inches(4.1), Inches(9.0), Inches(0.3),
                "Analysis Output:", font_size=15, bold=True, color=DARK_BLUE)

    output_code = """\
{ "arrays_needing_copy": ["a"],
  "war_dependencies": [{"array": "a", "read_idx": "j", "write_idx": "i"}] }"""
    add_code_box(slide, Inches(0.5), Inches(4.4), Inches(7.0), Inches(0.6), output_code, font_size=11)

    # Solution
    sol_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.5), Inches(5.15), Inches(9.0), Inches(0.45))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    sol_box.line.color.rgb = ACCENT_GREEN
    sol_box.line.width = Pt(1)
    sol_box.adjustments[0] = 0.1
    tf = sol_box.text_frame
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Injected into prompt: "
    run.font.name = FONT_BODY
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = ACCENT_GREEN
    run2 = p.add_run()
    run2.text = '"Use read-only copy:  a_copy = a.clone()  before the parallel region"'
    run2.font.name = FONT_CODE
    run2.font.size = Pt(13)
    run2.font.color.rgb = DARK_TEXT


def slide_07_parallel_reduction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Parallelization & Reduction Analysis")

    # GEMM example
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.3),
                "GEMM: C[i][j] += alpha * A[i][k] * B[k][j]",
                font_size=16, bold=True, color=MED_BLUE)

    # Parallel dims table
    rows = [
        ["Dimension", "Role", "Triton Mapping"],
        ["i", "Parallel (outer)", "program_id(0) -> block of rows"],
        ["j", "Parallel (inner)", "program_id(1) -> block of columns"],
        ["k", "Reduction (sequential)", "In-kernel for loop with accumulation"],
    ]
    add_simple_table(slide, Inches(0.5), Inches(1.5), Inches(9.0),
                     [1.2, 2.5, 5.3], rows)

    # Reduction classification
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(9.0), Inches(0.3),
                "Reduction Classification:", font_size=15, bold=True, color=DARK_BLUE)

    add_bullet_list(slide, Inches(0.5), Inches(3.5), Inches(9.0), Inches(1.5), [
        'Pattern: C[i][j] += A[i][k] * B[k][j]  ->  classified as "dot product"',
        "Triton strategy: 2D tiling + k-loop accumulation via tl.sum(A * B)",
        "Both i and j are safe to parallelize (no cross-iteration deps on C)",
    ], font_size=13)

    # Prompt injection
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(9.0), Inches(0.3),
                "Injected into Prompt:", font_size=15, bold=True, color=DARK_BLUE)

    prompt_code = """\
## Parallelization Analysis
- Parallelize `i`, sequential `j`: VALID
- Parallelize `j`, sequential `i`: VALID
- Parallelize both `i` and `j`: VALID (2D grid)
## Reduction: k-dimension (dot product pattern)"""
    add_code_box(slide, Inches(0.5), Inches(4.9), Inches(7.0), Inches(1.1), prompt_code, font_size=11)


def slide_08_prompt(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Analysis -> LLM Prompt")

    add_textbox(slide, Inches(0.5), Inches(0.9), Inches(9.0), Inches(0.35),
                "Real prompt structure for GEMM (from raw_responses/gemm/attempt1.txt):",
                font_size=13, color=MED_GRAY)

    # Prompt sections as boxes
    sections = [
        ("1. C Source Code", "Full kernel with #pragma scop annotations",
         DARK_BLUE, Inches(1.3)),
        ("2. Parallelization Analysis", 'Which dims are parallel, which are reduction\n'
         '"Parallelize i, sequential j: VALID"',
         MED_BLUE, Inches(1.85)),
        ("3. Array Information", "A: read-only | B: read-only | C: read-write",
         LIGHT_BLUE, Inches(2.55)),
        ("4. Dimension Parameters", "NI=60, NJ=70, NK=80 (compile-time -> runtime)",
         MED_BLUE, Inches(3.1)),
        ("5. Function Signature", 'def gemm_triton(A, B, C, alpha, beta, NI, NJ, NK)',
         LIGHT_BLUE, Inches(3.65)),
        ("6. Triton Compilation Rules", "Critical patterns: tl.arange outside loops,\n"
         "masked loads, constexpr block sizes",
         ACCENT_ORANGE, Inches(4.2)),
    ]

    for title, desc, color, y in sections:
        add_flow_box(slide, Inches(0.5), y, Inches(2.5), Inches(0.45), title,
                     fill_color=color, font_size=10)
        add_textbox(slide, Inches(3.2), y + Inches(0.02), Inches(6.5), Inches(0.45),
                    desc, font_size=11, color=DARK_TEXT)

    # Key insight
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.5), Inches(5.0), Inches(9.0), Inches(0.5))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
    insight_box.line.color.rgb = ACCENT_ORANGE
    insight_box.line.width = Pt(1)
    insight_box.adjustments[0] = 0.1
    tf = insight_box.text_frame
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key: "
    run.font.bold = True
    run.font.size = Pt(14)
    run.font.name = FONT_BODY
    run.font.color.rgb = ACCENT_ORANGE
    run2 = p.add_run()
    run2.text = "Analysis data is embedded as natural language guidance, not hidden metadata"
    run2.font.size = Pt(14)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


def slide_09_gemm_triton(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Generated Triton Code -- GEMM")

    # C code (left)
    add_textbox(slide, Inches(0.3), Inches(0.9), Inches(4.0), Inches(0.3),
                "C Original (6 lines):", font_size=12, bold=True, color=DARK_BLUE)
    c_code = """\
for (i = 0; i < NI; i++) {
  for (j = 0; j < NJ; j++)
    C[i][j] *= beta;
  for (k = 0; k < NK; k++)
    for (j = 0; j < NJ; j++)
      C[i][j] += alpha*A[i][k]*B[k][j];
}"""
    add_code_box(slide, Inches(0.3), Inches(1.2), Inches(4.5), Inches(1.55), c_code, font_size=10)

    # Triton code (right)
    add_textbox(slide, Inches(5.1), Inches(0.9), Inches(4.8), Inches(0.3),
                "Generated Triton Kernel:", font_size=12, bold=True, color=ACCENT_GREEN)
    triton_code = """\
@triton.jit
def gemm_kernel(A, B, C, alpha, beta,
                NI, NJ, NK,
                BLOCK_I: tl.constexpr,
                BLOCK_J: tl.constexpr):
  pid_i = tl.program_id(0)
  pid_j = tl.program_id(1)
  i_off = pid_i*BLOCK_I + tl.arange(0,BLOCK_I)
  j_off = pid_j*BLOCK_J + tl.arange(0,BLOCK_J)
  # C[i][j] *= beta
  c_idx = i_off[:,None]*NJ + j_off[None,:]
  mask = (i_off[:,None]<NI) & (j_off[None,:]<NJ)
  c = tl.load(C+c_idx, mask=mask) * beta
  # k-loop accumulation
  for k in range(NK):
    a = tl.load(A + i_off*NK+k, mask=i_off<NI)
    b = tl.load(B + k*NJ+j_off, mask=j_off<NJ)
    c += alpha * a[:,None] * b[None,:]
  tl.store(C+c_idx, c, mask=mask)"""
    add_code_box(slide, Inches(5.1), Inches(1.2), Inches(4.7), Inches(3.5), triton_code, font_size=9)

    # Highlights
    add_textbox(slide, Inches(0.3), Inches(3.0), Inches(4.5), Inches(0.3),
                "Key Triton Patterns:", font_size=13, bold=True, color=MED_BLUE)
    add_bullet_list(slide, Inches(0.3), Inches(3.3), Inches(4.5), Inches(2.0), [
        "2D tiling: pid_i, pid_j -> block of C",
        "Masked loads: boundary safety",
        "k-loop: sequential accumulation",
        "Broadcasting: a[:,None] * b[None,:]",
        "Scalar broadcast: alpha, beta",
    ], font_size=12)

    # Result badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.3), Inches(5.0), Inches(3.5), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_GREEN
    badge.line.fill.background()
    badge.adjustments[0] = 0.15
    tf = badge.text_frame
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "First-try pass  |  3.00x speedup"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def slide_10_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Results Summary")

    # Correctness stats
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(4.0), Inches(0.3),
                "Correctness: 25/30 (83.3%)",
                font_size=18, bold=True, color=ACCENT_GREEN)

    add_bullet_list(slide, Inches(0.5), Inches(1.4), Inches(4.0), Inches(0.8), [
        "11 first-try passes",
        "14 after retry (up to 5 attempts)",
        "vs TSVC: 150/151 (99.3%)",
    ], font_size=13)

    # Performance stats
    add_textbox(slide, Inches(5.0), Inches(1.0), Inches(4.5), Inches(0.3),
                "Performance",
                font_size=18, bold=True, color=MED_BLUE)

    add_bullet_list(slide, Inches(5.0), Inches(1.4), Inches(4.5), Inches(0.8), [
        "Median speedup: 1.76x",
        "Mean speedup: 2.05x",
        "GPU wins (>1x): 18/25 (72%)",
    ], font_size=13)

    # Top 5 table
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(9.0), Inches(0.3),
                "Top Speedups:", font_size=15, bold=True, color=DARK_BLUE)

    top_rows = [
        ["Kernel", "Speedup", "C Time (ms)", "Triton Time (ms)"],
        ["covariance", "5.40x", "0.340", "0.063"],
        ["deriche", "5.28x", "0.423", "0.080"],
        ["2mm", "3.96x", "0.573", "0.145"],
        ["trmm", "3.26x", "0.200", "0.061"],
        ["lu", "3.12x", "0.224", "0.072"],
        ["syrk", "3.08x", "0.178", "0.058"],
        ["gemm", "3.00x", "0.174", "0.058"],
    ]
    add_simple_table(slide, Inches(0.5), Inches(2.85), Inches(9.0),
                     [2.2, 2.2, 2.3, 2.3], top_rows)

    # Slowdowns note
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(9.0), Inches(0.3),
                "Slowdowns: nussinov (0.04x), trisolv (0.04x) -- sequential algorithms on GPU",
                font_size=12, color=ACCENT_RED)


def slide_11_failures(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Failure Analysis & Insights")

    # Failing kernels table
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.3),
                "5 Failing Kernels (all fundamentally sequential):",
                font_size=15, bold=True, color=ACCENT_RED)

    fail_rows = [
        ["Kernel", "Root Cause"],
        ["doitgen", "3D tensor contraction with sequential accumulation"],
        ["durbin", "Levinson-Durbin recursion: each step depends on all previous"],
        ["gramschmidt", "Sequential orthogonalization: each vector depends on prior"],
        ["ludcmp", "LU decomposition: sequential pivot-dependent updates"],
        ["seidel_2d", "Gauss-Seidel: reads current-iteration neighbors"],
    ]
    add_simple_table(slide, Inches(0.5), Inches(1.35), Inches(9.0),
                     [2.0, 7.0], fail_rows)

    # 7 slowdown kernels
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(9.0), Inches(0.3),
                "7 Slowdown Kernels (<1x): sequential algorithms that don't benefit from GPU",
                font_size=13, bold=True, color=ACCENT_ORANGE)

    # Key insight
    y_insight = Inches(4.3)
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.5), y_insight, Inches(9.0), Inches(1.2))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    insight_box.line.color.rgb = MED_BLUE
    insight_box.line.width = Pt(1.5)
    insight_box.adjustments[0] = 0.05

    tf = insight_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Insight: "
    run.font.bold = True
    run.font.size = Pt(15)
    run.font.name = FONT_BODY
    run.font.color.rgb = MED_BLUE

    run2 = p.add_run()
    run2.text = ("Static analysis correctly identifies parallelizable vs sequential kernels. "
                 "The LLM leverages this guidance effectively, but cannot overcome "
                 "fundamental algorithmic limits -- sequential dependencies are a property "
                 "of the algorithm, not a failure of the pipeline.")
    run2.font.size = Pt(14)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


# ════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    slide_01_title(prs)
    slide_02_architecture(prs)
    slide_03_pet(prs)
    slide_04_llvm(prs)
    slide_05_fallback(prs)
    slide_06_war(prs)
    slide_07_parallel_reduction(prs)
    slide_08_prompt(prs)
    slide_09_gemm_triton(prs)
    slide_10_results(prs)
    slide_11_failures(prs)

    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            "polybench_pipeline_slides.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
